#!/usr/bin/env python3
"""
ТЦ Бродвей — повна презентація (15 слайдів) у стилі МПЗ Сакура
Формат: 10" × 5.625"

Структура:
  Слайди 1-10  — основний контент
  Слайди 11-14 — фотогалерея
  Слайд  15    — контакти
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ═══════════════════════════════════════════════════════════════
# КОЛЬОРИ
# ═══════════════════════════════════════════════════════════════
BG     = RGBColor(0x0C, 0x18, 0x28)
HBAR   = RGBColor(0x0F, 0x1B, 0x2D)
CARD   = RGBColor(0x1A, 0x35, 0x50)
CARD2  = RGBColor(0x14, 0x29, 0x42)
GOLD_S = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_T = RGBColor(0xE8, 0xC9, 0x7A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LABEL  = RGBColor(0xD6, 0xE0, 0xEB)
MUTED  = RGBColor(0x8A, 0x9B, 0xB0)
DARK_T = RGBColor(0x0F, 0x1B, 0x2D)
DARK_S = RGBColor(0x1A, 0x35, 0x50)
GREEN  = RGBColor(0x2D, 0x7D, 0x57)
GREEN2 = RGBColor(0x1A, 0x4A, 0x32)
RED    = RGBColor(0xC0, 0x39, 0x2B)
RED2   = RGBColor(0x5A, 0x18, 0x12)
BLUE2  = RGBColor(0x12, 0x2A, 0x4A)
AMBER2 = RGBColor(0x4A, 0x35, 0x10)
HL     = RGBColor(0x1E, 0x42, 0x62)   # highlighted row

# ═══════════════════════════════════════════════════════════════
# РОЗМІРИ
# ═══════════════════════════════════════════════════════════════
W = Inches(10.0)
H = Inches(5.625)

P = {
    "main":  "/sessions/zealous-beautiful-volta/mnt/Claude Code/data/objects/БЦ Торговое/photos/main foto.jpeg",
    "p1":    "/sessions/zealous-beautiful-volta/mnt/Claude Code/data/objects/БЦ Торговое/photos/1.jpeg",
    "p2":    "/sessions/zealous-beautiful-volta/mnt/Claude Code/data/objects/БЦ Торговое/photos/2.jpeg",
    "p3":    "/sessions/zealous-beautiful-volta/mnt/Claude Code/data/objects/БЦ Торговое/photos/3.jpeg",
    "p4":    "/sessions/zealous-beautiful-volta/mnt/Claude Code/data/objects/БЦ Торговое/photos/4.jpeg",
    "p5":    "/sessions/zealous-beautiful-volta/mnt/Claude Code/data/objects/БЦ Торговое/photos/5.jpeg",
    "geo":   "/sessions/zealous-beautiful-volta/mnt/Claude Code/data/objects/БЦ Торговое/photos/geo.png",
}

NB = "\u00A0"   # non-breaking space

# ═══════════════════════════════════════════════════════════════
# БАЗОВІ УТИЛІТИ
# ═══════════════════════════════════════════════════════════════

def R(slide, x, y, w, h, fill, border=None, bw=0.5):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def T(slide, text, x, y, w, h, size, bold=False, color=WHITE,
      align=PP_ALIGN.LEFT, v="t", wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    bp = tf._txBody.find(qn("a:bodyPr"))
    if bp is None:
        bp = etree.SubElement(tf._txBody, qn("a:bodyPr"))
    bp.set("anchor", v)
    bp.set("wrap", "none" if not wrap else "square")
    for a in ("lIns", "rIns", "tIns", "bIns"):
        bp.set(a, "0")
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Nunito Sans"
    return tb


def pic(slide, path, x, y, w, h):
    """Вставляє фото з fallback-заглушкою."""
    try:
        slide.shapes.add_picture(path, x, y, w, h)
    except Exception:
        R(slide, x, y, w, h, CARD)
        T(slide, "[ Фото ]", x, y, w, h, 12, color=MUTED,
          align=PP_ALIGN.CENTER, v="ctr")


def bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG


# ═══════════════════════════════════════════════════════════════
# КОМПОНЕНТИ
# ═══════════════════════════════════════════════════════════════

def gold_vbar(slide):
    R(slide, 0, 0, Inches(0.12), H, GOLD_S)


def header_bar(slide, title, tag="", fill=HBAR):
    R(slide, 0, 0, W, Inches(0.82), fill)
    T(slide, title, Inches(0.28), 0, W * 0.72, Inches(0.82),
      17, bold=True, color=WHITE, v="ctr")
    if tag:
        T(slide, tag, Inches(7.3), 0, Inches(2.5), Inches(0.82),
          9, color=GOLD_T, align=PP_ALIGN.RIGHT, v="ctr")


def footer_bar(slide, text, fill=HBAR):
    R(slide, 0, H - Inches(0.30), W, Inches(0.30), fill)
    T(slide, text,
      Inches(0.40), H - Inches(0.30), W - Inches(0.5), Inches(0.30),
      7, color=MUTED, v="ctr", wrap=False)


def kpi_card_cover(slide, x, y, w, h, value, label, val_sz=20, lbl_sz=10):
    R(slide, x, y, w, h, CARD)
    R(slide, x, y, w, Inches(0.055), GOLD_S)
    T(slide, value, x, y + Inches(0.10), w, h * 0.55,
      val_sz, bold=True, color=GOLD_T,
      align=PP_ALIGN.CENTER, v="ctr", wrap=False)
    T(slide, label, x, y + h * 0.66, w, h * 0.30,
      lbl_sz, color=LABEL, align=PP_ALIGN.CENTER, v="ctr", wrap=False)


def kpi_card_big(slide, x, y, w, h, value, label):
    R(slide, x, y, w, h, CARD)
    R(slide, x, y, w, Inches(0.07), GOLD_S)
    T(slide, value, x, y + Inches(0.15), w, h * 0.52,
      26, bold=True, color=GOLD_T, align=PP_ALIGN.CENTER, v="ctr")
    T(slide, label, x, y + h * 0.72, w, h * 0.28,
      9, color=LABEL, align=PP_ALIGN.CENTER, v="ctr", wrap=True)


def info_card_white(slide, x, y, w, h, letter, title, body):
    R(slide, x, y, w, h, WHITE)
    ico = Inches(0.52)
    R(slide, x + Inches(0.10), y + (h - ico) / 2, ico, ico, CARD)
    T(slide, letter, x + Inches(0.10), y + (h - ico) / 2, ico, ico,
      13, bold=True, color=GOLD_T, align=PP_ALIGN.CENTER, v="ctr")
    T(slide, title, x + Inches(0.73), y + Inches(0.08),
      w - Inches(0.83), Inches(0.28), 9, bold=True, color=DARK_S)
    T(slide, body, x + Inches(0.73), y + Inches(0.37),
      w - Inches(0.83), h - Inches(0.44), 9, color=DARK_T, wrap=True)


def stat_row(slide, x, y, w, h, value, label, val_color=GOLD_T):
    R(slide, x, y, w, h, CARD)
    R(slide, x, y, Inches(0.06), h, GOLD_S)
    T(slide, value, x + Inches(0.18), y, w * 0.38, h,
      12, bold=True, color=val_color, v="ctr")
    T(slide, label, x + w * 0.40, y, w * 0.57, h,
      9, color=LABEL, v="ctr", wrap=True)


def section_label(slide, text):
    """Маленький тег секції в правому верхньому куті."""
    T(slide, text, Inches(7.3), 0, Inches(2.5), Inches(0.82),
      8, color=GOLD_T, align=PP_ALIGN.RIGHT, v="ctr")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 1 — ОБКЛАДИНКА
# ═══════════════════════════════════════════════════════════════

def slide1_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    R(s, Inches(6.5), 0, Inches(3.5), H, CARD)
    pic(s, P["main"], Inches(6.5), 0, Inches(3.5), H)
    gold_vbar(s)
    R(s, Inches(6.50) - Pt(1), 0, Pt(1.5), H, GOLD_S)

    R(s, Inches(0.50), Inches(0.46), Inches(2.7), Inches(0.30), GOLD_S)
    T(s, "ІНВЕСТИЦІЙНА ПРЕЗЕНТАЦІЯ",
      Inches(0.52), Inches(0.46), Inches(2.66), Inches(0.30),
      7, bold=True, color=DARK_T, v="ctr", wrap=False)

    T(s, "ТЦ БРОДВЕЙ",
      Inches(0.40), Inches(0.90), Inches(5.8), Inches(1.15),
      52, bold=True, color=WHITE, v="ctr")

    T(s, f"вул.{NB}Таращанська,{NB}191-А  ·  Біла{NB}Церква",
      Inches(0.40), Inches(2.05), Inches(5.8), Inches(0.50),
      18, color=GOLD_T, v="ctr", wrap=False)

    R(s, Inches(0.40), Inches(2.65), Inches(5.8), Pt(1.5), GOLD_S)

    T(s, "Районний торговий центр  |  Клас B  |  Реконструкція 2018",
      Inches(0.40), Inches(2.80), Inches(5.8), Inches(0.82),
      12, color=LABEL, v="ctr")

    kpis = [
        (f"$3{NB}000{NB}000", "Ціна пропозиції"),
        ("15,27%",            "Cap Rate"),
        ("6,5 р.",            "Термін окупності"),
    ]
    ML = Inches(0.40); MR = Inches(0.40); GAP = Inches(0.10)
    kw = (Inches(6.5) - ML - MR - GAP * 2) / 3
    ky = Inches(3.75);  kh = H - ky - Inches(0.37)
    for i, (v, l) in enumerate(kpis):
        kpi_card_cover(s, ML + i * (kw + GAP), ky, kw, kh, v, l,
                       val_sz=16, lbl_sz=10)

    footer_bar(s, f"БІЛА{NB}ЦЕРКВА, КИЇВСЬКА{NB}ОБЛАСТЬ  |  208{NB}000 МЕШ.  |  80 КМ ВІД КИЄВА  |  2026")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 2 — ОГЛЯД ОБ'ЄКТА
# ═══════════════════════════════════════════════════════════════

def slide2_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); header_bar(s, "ОГЛЯД ОБ'ЄКТА", "ТЦ «Бродвей»"); gold_vbar(s)

    cards = [
        ("TC",  "Тип об'єкта",
         "Районний ТЦ: рітейл-галерея + якірний супермаркет Варус (з 01.05.2026)"),
        ("Pin", "Локація",
         "Район Заріччя, вул. Таращанська 191-А, Біла Церква (208 000 мешк., 80 км від Києва)"),
        ("M2",  "Площа та будівля",
         "GBA: 3 486 м²  |  GLA: 3 258 м²  |  1 поверх  |  Цегла  |  Клас B"),
        ("OK",  "Стан та документи",
         "Реконструкція 2018, хороший стан. Право власності зареєстровано в ЄДРПР"),
    ]
    cw = Inches(4.40); ch = Inches(0.88); cx = Inches(0.30)
    y0 = Inches(0.98); gap = Inches(0.10)
    for i, (le, ti, bo) in enumerate(cards):
        info_card_white(s, cx, y0 + i * (ch + gap), cw, ch, le, ti, bo)

    stats = [
        ("13",      "орендарів — стабільне заповнення"),
        ("93,5%",   "зайнятість площ"),
        ("Варус",   "якірний орендар (з 01.05.2026)"),
        ("0,47 га", "земельна ділянка (оренда до 2027)"),
        (f"$861/м²","ціна за 1 м² GBA"),
    ]
    rw = Inches(4.72); rh = Inches(0.70); rx = Inches(5.00)
    ry0 = Inches(0.98); rg = Inches(0.085)
    for i, (v, l) in enumerate(stats):
        stat_row(s, rx, ry0 + i * (rh + rg), rw, rh, v, l)

    footer_bar(s, f"GBA: 3{NB}486 м²  |  GLA: 3{NB}258 м²  |  Поверхів: 1  |  Клас B  |  Реконструкція: 2018  |  Ціна: $3{NB}000{NB}000")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 3 — ЛОКАЦІЯ ТА МІСТО
# ═══════════════════════════════════════════════════════════════

def slide3_location(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); header_bar(s, "ЛОКАЦІЯ ТА МІСТО", "ТЦ «Бродвей»"); gold_vbar(s)

    # Карта (ліва частина)
    map_w = Inches(4.55); map_x = Inches(0.30)
    map_y = Inches(0.96); map_h = H - map_y - Inches(0.34)
    pic(s, P["geo"], map_x, map_y, map_w, map_h)
    R(s, map_x + map_w - Pt(1), map_y, Pt(1), map_h, GOLD_S)

    # Права колонка — характеристики міста та локації
    stats = [
        ("208 000",  "мешканців — 32-е місто України"),
        ("80 км",    "від Києва (траса E95/M05)"),
        ("Заріччя",  "житловий район — стабільний трафік"),
        ("2 школи",  "поруч (№13, №22) + Укрпошта"),
        ("Паркінг",  "власне відкрите паркування перед входом"),
        ("E95/M05",  "міжнародна траса поблизу міста"),
    ]
    rw = Inches(4.78); rh = Inches(0.64); rx = Inches(5.00)
    ry0 = Inches(0.96); rg = Inches(0.073)
    for i, (v, l) in enumerate(stats):
        stat_row(s, rx, ry0 + i * (rh + rg), rw, rh, v, l)

    footer_bar(s, "Біла Церква — найбільш промислово розвинене місто Київської обл.  |  57+ підприємств  |  Розвинений сектор послуг")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 4 — ТЕХНІЧНІ ХАРАКТЕРИСТИКИ
# ═══════════════════════════════════════════════════════════════

def slide4_technical(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); header_bar(s, "ТЕХНІЧНІ ХАРАКТЕРИСТИКИ", "ТЦ «Бродвей»"); gold_vbar(s)

    cards = [
        ("GBA", "Загальна площа (GBA)",   "3 485,9 м² — підтверджено Витягом з ЄДРПР"),
        ("GLA", "Орендована площа (GLA)", "3 258,33 м² — активно здається в оренду"),
        ("1F",  "Поверховість",           "1 поверх. Одноповерхова будівля з зовнішньою аркадою"),
        ("TC",  "Формат та оздоблення",   "Торгова галерея + якірний супермаркет. Плитка, LED-освітлення, підвісні стелі"),
    ]
    cw = Inches(4.40); ch = Inches(0.88); cx = Inches(0.30)
    y0 = Inches(0.98); gap = Inches(0.10)
    for i, (le, ti, bo) in enumerate(cards):
        info_card_white(s, cx, y0 + i * (ch + gap), cw, ch, le, ti, bo)

    stats = [
        ("Цегла",  "тип конструкції будівлі"),
        ("2018",   "рік реконструкції"),
        ("Клас B", "клас торгової нерухомості"),
        ("Плитка", "підлога (торгова зона + зовн. аркада)"),
        ("24/7",   "охорона та відеоспостереження"),
    ]
    rw = Inches(4.72); rh = Inches(0.70); rx = Inches(5.00)
    ry0 = Inches(0.98); rg = Inches(0.085)
    for i, (v, l) in enumerate(stats):
        stat_row(s, rx, ry0 + i * (rh + rg), rw, rh, v, l)

    footer_bar(s, f"GBA: 3{NB}485,9 м²  |  GLA: 3{NB}258,33 м²  |  Поверхів: 1  |  Цегла  |  Реконструкція 2018  |  Клас B")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 5 — ЗЕМЕЛЬНА ДІЛЯНКА ТА ПРАВОВИЙ СТАТУС
# ═══════════════════════════════════════════════════════════════

def slide5_legal(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); header_bar(s, "ЗЕМЕЛЬНА ДІЛЯНКА ТА ПРАВОВИЙ СТАТУС", "ТЦ «Бродвей»"); gold_vbar(s)

    # Ліва колонка — документи (білі картки)
    cards = [
        ("Doc", "Право власності",
         "Зареєстровано в ЄДРПР. Витяг на нежитлове приміщення + Договір купівлі-продажу"),
        ("Map", "Земельна ділянка",
         f"Кадастр: 3220489500:02:021:0736  |  0,4679 га  |  Оренда у міської ради"),
        ("Cal", "Строк оренди землі",
         "5 років: 11.08.2022 – 11.08.2027. Автоматичне продовження. Орендна плата 529 159 грн/рік"),
        ("OK",  "Дозвільна документація",
         "Дозвіл на будівництво + Договір про реконструкцію + Акт + Сертифікат"),
    ]
    cw = Inches(4.40); ch = Inches(0.88); cx = Inches(0.30)
    y0 = Inches(0.98); gap = Inches(0.10)
    for i, (le, ti, bo) in enumerate(cards):
        info_card_white(s, cx, y0 + i * (ch + gap), cw, ch, le, ti, bo)

    # Права колонка — статус
    stats = [
        ("ЄДРПР",  "право власності зареєстровано"),
        ("0,47 га","площа земельної ділянки"),
        ("03.07",  "цільове призначення — будівлі торгівлі"),
        ("2027",   "строк оренди землі (з автопродовженням)"),
        ("Повний", "пакет правовстановлюючих документів"),
    ]
    rw = Inches(4.72); rh = Inches(0.70); rx = Inches(5.00)
    ry0 = Inches(0.98); rg = Inches(0.085)
    for i, (v, l) in enumerate(stats):
        stat_row(s, rx, ry0 + i * (rh + rg), rw, rh, v, l)

    footer_bar(s, "Власник: фізична особа-підприємець  |  Земля: комунальна власність (Білоцерківська міська рада)  |  Договорів з комунальниками: 2")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 6 — ОРЕНДАРІ ТА ДОХОДИ
# ═══════════════════════════════════════════════════════════════

def slide6_tenants(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); header_bar(s, "ОРЕНДАРІ ТА ДОХОДИ", "ТЦ «Бродвей»"); gold_vbar(s)

    # ── Ліва частина: таблиця орендарів ───────────────────────
    tx = Inches(0.28); ty = Inches(0.96)
    tw = Inches(5.55); col = [Inches(1.55), Inches(1.45), Inches(1.75), Inches(0.80)]
    hdr = ["Орендар",  "Площа, м²",  "Оренда, грн/міс", "Частка"]
    rows = [
        ["Варус (з 01.05.2026)", "1 750,76", "1 417 000",  "81,9%"],
        ["Аптека",               "80,0",     "н/д",        "~2,5%"],
        ["Мобільний оператор",   "40,0",     "н/д",        "~1,2%"],
        ["Товари для дому",      "120,0",    "н/д",        "~3,7%"],
        ["Одяг та аксесуари",    "150,0",    "н/д",        "~4,6%"],
        ["Інші 8 орендарів",     "1 117,57", "312 940",    "6,1%"],
        ["РАЗОМ (13 орендарів)", "3 258,33", "1 729 940",  "100%"],
    ]
    rh = Inches(0.49)
    # Заголовок таблиці
    cx = tx
    for ci, (cw, ht) in enumerate(zip(col, hdr)):
        R(s, cx, ty, cw, rh, GOLD_S)
        T(s, ht, cx + Inches(0.04), ty, cw - Inches(0.06), rh,
          8, bold=True, color=DARK_T, v="ctr", wrap=False)
        cx += cw
    # Рядки даних
    for ri, row in enumerate(rows):
        ry = ty + rh + ri * (rh - Inches(0.01))
        is_last = (ri == len(rows) - 1)
        is_varus = (ri == 0)
        fill = HL if is_varus else (CARD if ri % 2 == 0 else BG)
        if is_last:
            fill = RGBColor(0x0A, 0x1E, 0x35)
        cx = tx
        for ci, (cw, cell) in enumerate(zip(col, row)):
            R(s, cx, ry, cw, rh - Inches(0.01), fill)
            fc = GOLD_T if (is_varus or is_last) else (WHITE if ci > 0 else LABEL)
            T(s, cell, cx + Inches(0.04), ry, cw - Inches(0.06), rh - Inches(0.01),
              8 if ci > 0 else 8, bold=(is_varus or is_last), color=fc,
              v="ctr", wrap=False,
              align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
            cx += cw

    # ── Права частина: KPI орендарів ──────────────────────────
    rx2 = Inches(6.05)
    stats = [
        ("13",         "орендарів"),
        ("93,5%",      "зайнятість площ"),
        (f"$39{NB}769","дохід на місяць (USD)"),
        ("81,9%",      "дохід від якірного орендаря"),
        (f"$12,21/м²", "середня ставка оренди"),
    ]
    rw2 = Inches(3.67); rh2 = Inches(0.70)
    ry0 = Inches(0.96); rg2 = Inches(0.085)
    for i, (v, l) in enumerate(stats):
        stat_row(s, rx2, ry0 + i * (rh2 + rg2), rw2, rh2, v, l)

    footer_bar(s, f"Місячний дохід: 1{NB}729{NB}940 грн / $39{NB}769  |  EGI: $477{NB}225/рік  |  Вакансія: 6,5%  |  Якір: Варус (1{NB}751 м²)")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 7 — ФІНАНСОВА МОДЕЛЬ
# ═══════════════════════════════════════════════════════════════

def slide7_finance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); header_bar(s, "ФІНАНСОВА МОДЕЛЬ", "ТЦ «Бродвей»", fill=CARD); gold_vbar(s)

    kpis = [
        (f"$458{NB}131", "NOI (чистий операц. дохід)"),
        (f"$477{NB}225", "EGI (річний дохід оренди)"),
        ("15,27%",       "Cap Rate"),
        ("6,5 р.",       "Термін окупності"),
    ]
    cw = Inches(2.15); ch = Inches(1.35)
    gap = Inches(0.115); sx = Inches(0.30)
    for i, (v, l) in enumerate(kpis):
        kpi_card_big(s, sx + i * (cw + gap), Inches(0.96), cw, ch, v, l)

    # ── Доходи та витрати (ліво) ──────────────────────────────
    bx = Inches(0.30); by = Inches(2.52)
    bw = Inches(5.50); bh = Inches(0.50); bg2 = Inches(0.08)

    T(s, "СТРУКТУРА ДОХОДІВ ТА ВИТРАТ",
      bx, by, bw, Inches(0.26), 8, bold=True, color=GOLD_T)
    by += Inches(0.28)

    MAX = 477225.0
    bars = [
        ("EGI (ефект. дохід)",     f"$477{NB}225", 477225/MAX, GOLD_S),
        ("Оренда землі",           "$12{NB}165",   12165/MAX,  RGBColor(0x4A,0x7A,0xAF)),
        ("Податок на майно",       "$6{NB}929",    6929/MAX,   RGBColor(0x5A,0x50,0x8E)),
        ("NOI (чистий дохід)",     f"$458{NB}131", 458131/MAX, GREEN),
    ]
    LW = Inches(1.30); VW = Inches(0.85)
    MAXBW = bw - LW - VW - Inches(0.10)
    for lbl, vstr, ratio, color in bars:
        R(s, bx, by, bw, bh, CARD)
        T(s, lbl, bx + Inches(0.08), by, LW, bh, 8, color=LABEL, v="ctr")
        barw = max(MAXBW * ratio, Inches(0.06))
        R(s, bx + LW, by + Inches(0.07), barw, bh - Inches(0.14), color)
        T(s, vstr, bx + LW + MAXBW + Inches(0.05), by, VW, bh,
          8, bold=True, color=GOLD_T, align=PP_ALIGN.RIGHT, v="ctr", wrap=False)
        by += bh + bg2

    # ── Сценарії (право) ──────────────────────────────────────
    tx = Inches(6.10); ty = Inches(2.52)
    tw2 = Inches(3.65); th = H - ty - Inches(0.34)

    T(s, "СЦЕНАРІЇ ОЦІНКИ (за NOI $458 131/рік)",
      tx, ty, tw2, Inches(0.26), 8, bold=True, color=GOLD_T)

    rows = [
        ("Окупність", "Вартість",  "Cap Rate", "Δ від ціни"),
        ("4,5 р.",    "$2 062К",   "22,2%",    "−31,3%"),
        ("5,5 р.",    "$2 520К",   "18,2%",    "−16,0%"),
        ("6,5 р. ✦",  "$2 978К",   "15,4%",    "−0,7%"),
        ("7,5 р.",    "$3 436К",   "13,3%",    "+14,5%"),
        ("8,5 р.",    "$3 894К",   "11,8%",    "+29,8%"),
    ]
    cws = [tw2*0.24, tw2*0.27, tw2*0.22, tw2*0.27]
    trh = (th - Inches(0.28)) / len(rows)
    ty2 = ty + Inches(0.28)
    for ri, row in enumerate(rows):
        is_hdr  = ri == 0; is_base = ri == 3
        rf = GOLD_S if is_hdr else (HL if is_base else (CARD if ri%2==0 else BG))
        cx = tx
        for ci, cell in enumerate(row):
            cw2 = cws[ci]
            R(s, cx, ty2, cw2, trh, rf)
            if ci < 3:
                R(s, cx+cw2-Pt(0.5), ty2, Pt(1), trh, BG)
            tc = DARK_T if is_hdr else (GOLD_T if (is_base or ci==0) else WHITE)
            T(s, cell, cx+Inches(0.04), ty2, cw2-Inches(0.06), trh,
              8 if is_hdr else 9, bold=(is_hdr or is_base), color=tc,
              align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT, v="ctr", wrap=False)
            cx += cw2
        ty2 += trh

    R(s, Inches(5.93), Inches(2.52), Pt(1), H-Inches(2.52)-Inches(0.34), CARD)
    footer_bar(s, f"EGI: $477{NB}225/рік  |  OPEX: $19{NB}094/рік (земля + податок)  |  NOI: $458{NB}131/рік  |  Ставка: $12,21/м²/міс  |  Зайнятість: 93,5%")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 8 — РИНКОВЕ ОТОЧЕННЯ
# ═══════════════════════════════════════════════════════════════

def slide8_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); header_bar(s, "РИНКОВЕ ОТОЧЕННЯ", "ТЦ «Бродвей»"); gold_vbar(s)

    # Ліва: конкуренти
    T(s, "КОНКУРЕНТИ В БІЛІЙ ЦЕРКВІ",
      Inches(0.30), Inches(0.96), Inches(4.55), Inches(0.28),
      8, bold=True, color=GOLD_T)

    comp_rows = [
        ("ТЦ", "Адреса",            "Площа",      "Позиція"),
        ("ТРК ВЕГА",    "вул. Героїв Небесної Сотні", "~13 000 м²",  "Головний ТРЦ"),
        ("ТРЦ Гермес",  "вул. Ярослава Мудрого",      "н/д",         "Центральний ТЦ"),
        ("Бульвар",     "бул. 50-річчя Перемоги",     "н/д",         "Торгово-офісний"),
        ("★ Бродвей",  "вул. Таращанська, 191-А",    "3 486 м²",    "Районний ТЦ"),
    ]
    cws2 = [Inches(0.90), Inches(1.65), Inches(1.05), Inches(0.95)]
    tx = Inches(0.30); ty = Inches(1.28); rh = Inches(0.52)
    for ri, row in enumerate(comp_rows):
        is_hdr = ri == 0; is_us = ri == 4
        rf = GOLD_S if is_hdr else (HL if is_us else (CARD if ri%2==0 else BG))
        cx = tx
        for ci, (cw2, cell) in enumerate(zip(cws2, row)):
            R(s, cx, ty, cw2, rh, rf)
            if ci < 3: R(s, cx+cw2-Pt(0.5), ty, Pt(1), rh, BG)
            tc = DARK_T if is_hdr else (GOLD_T if is_us else (LABEL if ci==0 else WHITE))
            T(s, cell, cx+Inches(0.04), ty, cw2-Inches(0.06), rh,
              8, bold=(is_hdr or is_us), color=tc, v="ctr", wrap=False,
              align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)
            cx += cw2
        ty += rh

    # Позиція Бродвей
    T(s, "ПОЗИЦІЯ ТЦ БРОДВЕЙ",
      Inches(0.30), Inches(4.00), Inches(4.55), Inches(0.26),
      8, bold=True, color=GOLD_T)
    R(s, Inches(0.30), Inches(4.28), Inches(4.55), Inches(0.90), CARD)
    T(s, "Нішевий районний ТЦ з фокусом на повсякденний попит мешканців "
         "Заріччя. Якір Варус формує стабільний трафік. Конкурує не з ВЕГА і "
         "Гермесом, а обслуговує локальну аудиторію у шаговій доступності.",
      Inches(0.38), Inches(4.30), Inches(4.40), Inches(0.86),
      9, color=LABEL, wrap=True)

    # Права: ринкові ставки та тренди
    T(s, "РИНОК ОРЕНДИ ТОРГОВИХ ПЛОЩ, БІЛА ЦЕРКВА",
      Inches(5.05), Inches(0.96), Inches(4.65), Inches(0.28),
      8, bold=True, color=GOLD_T)
    mstats = [
        ("$8–20/м²",  "ринкові ставки оренди ($/м²/міс)"),
        (f"$12,21/м²",f"поточна ставка Бродвей (в{NB}межах ринку)"),
        ("~10%",      "середній Cap Rate ринку торгової нерухомості"),
        ("15,27%",    "Cap Rate Бродвей — вище ринку на 5+ п.п."),
        ("Зростання", "попиту на районні ТЦ (шаговий сервіс)"),
    ]
    rw = Inches(4.65); rh = Inches(0.70); rx = Inches(5.05)
    ry = Inches(1.28); rg = Inches(0.085)
    for v, l in mstats:
        stat_row(s, rx, ry, rw, rh, v, l)
        ry += rh + rg

    footer_bar(s, "Ринок торг. нерухомості Київської обл.  |  Ставки: $8–20/м²/міс  |  Вакансія ринку: 8–12%  |  Бродвей: 6,5% — нижче ринку")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 9 — SWOT-АНАЛІЗ
# ═══════════════════════════════════════════════════════════════

def slide9_swot(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); header_bar(s, "SWOT-АНАЛІЗ", "ТЦ «Бродвей»"); gold_vbar(s)

    QY = Inches(0.90); QH = (H - QY - Inches(0.34)) / 2
    QW = (W - Inches(0.22)) / 2; GAP = Inches(0.06)

    quads = [
        # (x, y, fill_bg, label_bg, label, items)
        (Inches(0.22), QY,
         GREEN2, GREEN, "S  СИЛЬНІ СТОРОНИ",
         ["Реконструкція 2018 — мінімальний CAPEX на 5+ р.",
          "93,5% зайнятість — стабільний грошовий потік",
          "Якір Варус — довгостроковий договір, фіксована ставка",
          "Cap Rate 15,27% — значно вище ринку (~10%)",
          "Повний пакет правовстановлюючих документів"]),
        (Inches(0.22) + QW + GAP, QY,
         BLUE2, CARD, "W  СЛАБКІ СТОРОНИ",
         ["Варус генерує 81,9% доходу — ризик концентрації",
          "Земля в оренді у міськради (до 11.08.2027)",
          "1 поверх — обмежений потенціал розвитку",
          "Розташування: не центр міста (~2 км від центру)",
          "Відсутні дані про висоту стель та вент. системи"]),
        (Inches(0.22), QY + QH + GAP,
         AMBER2, RGBColor(0x7A,0x60,0x20), "O  МОЖЛИВОСТІ",
         ["Нові ЖК поблизу → зростання аудиторії",
          "Підвищення ставок при перезаключенні контрактів",
          "Продовження оренди землі на вигідних умовах",
          "Розвиток аркади: додаткові орендні площі",
          "Реінвестиція в об'єкт після виходу з позиції"]),
        (Inches(0.22) + QW + GAP, QY + QH + GAP,
         RED2, RED, "T  ЗАГРОЗИ",
         ["ТРК ВЕГА (13 000 м²) та ТРЦ Гермес — конкуренти",
          "Втрата якірного орендаря Варус — критичний ризик",
          "Військовий стан: невизначеність ринку нерухомості",
          "Зміна умов оренди землі після 2027 р.",
          "Зростання комунальних витрат та податкового тиску"]),
    ]

    for qx, qy, qfill, lfill, lbl, items in quads:
        R(s, qx, qy, QW, QH, qfill)
        # Label bar
        R(s, qx, qy, QW, Inches(0.30), lfill)
        T(s, lbl, qx + Inches(0.10), qy, QW - Inches(0.12), Inches(0.30),
          9, bold=True, color=WHITE, v="ctr", wrap=False)
        # Bullet items
        item_h = (QH - Inches(0.34)) / len(items)
        for ii, item in enumerate(items):
            iy = qy + Inches(0.32) + ii * item_h
            T(s, f"·  {item}", qx + Inches(0.10), iy, QW - Inches(0.15), item_h,
              8, color=LABEL, v="ctr", wrap=True)

    footer_bar(s, "Ключовий ризик: концентрація доходу на Варус (81,9%)  |  Ключова можливість: Cap Rate 15,27% вище ринку")


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 10 — ІНВЕСТИЦІЙНИЙ ВИСНОВОК
# ═══════════════════════════════════════════════════════════════

def slide10_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); header_bar(s, "ІНВЕСТИЦІЙНИЙ ВИСНОВОК", "ТЦ «Бродвей»"); gold_vbar(s)

    # 5 тез (ліворуч)
    T(s, "ІНВЕСТИЦІЙНІ ТЕЗИ",
      Inches(0.30), Inches(0.96), Inches(5.30), Inches(0.26),
      8, bold=True, color=GOLD_T)

    theses = [
        ("01", "Підтверджений грошовий потік",
         f"NOI $458{NB}131/рік підтверджено договорами. 93,5% площ зайнято."),
        ("02", "Cap Rate значно вище ринку",
         f"15,27% при ринковому рівні ~10% — надбавка 5+ п.п. за ризик."),
        ("03", "Якорний орендар знижує волатильність",
         "Варус — мережева структура, довгостроковий договір, фіксована ставка."),
        ("04", "Мінімальний CAPEX на горизонт 5+ р.",
         "Реконструкція 2018. Комерційне обладнання орендарів — за їх рахунок."),
        ("05", "Справедлива ціна відносно дохідності",
         f"При окупності 6,5 р. розрахункова вартість $2{NB}978К vs ціна $3{NB}000К (Δ −0,7%)."),
    ]
    ty = Inches(1.26); th = Inches(0.68); tg = Inches(0.07)
    for num, title, body in theses:
        # Номер
        R(s, Inches(0.30), ty, Inches(0.42), th, GOLD_S)
        T(s, num, Inches(0.30), ty, Inches(0.42), th,
          12, bold=True, color=DARK_T, align=PP_ALIGN.CENTER, v="ctr")
        # Контент
        R(s, Inches(0.74), ty, Inches(4.86), th, CARD)
        T(s, title, Inches(0.82), ty, Inches(4.70), Inches(0.26),
          9, bold=True, color=GOLD_T, v="ctr")
        T(s, body, Inches(0.82), ty + Inches(0.28), Inches(4.70), Inches(0.36),
          8, color=LABEL, v="ctr", wrap=True)
        ty += th + tg

    # Рекомендація (праворуч)
    rx = Inches(6.40); ry = Inches(0.96)
    rw = Inches(3.35); rh = H - ry - Inches(0.34)

    R(s, rx, ry, rw, rh, CARD)
    R(s, rx, ry, rw, Inches(0.07), GOLD_S)

    T(s, "РЕКОМЕНДАЦІЯ",
      rx, ry + Inches(0.12), rw, Inches(0.30),
      9, bold=True, color=GOLD_T, align=PP_ALIGN.CENTER, v="ctr")

    rec = ("ОБ'ЄКТ РЕКОМЕНДУЄТЬСЯ до інвестиційного розгляду як "
           "дохідна нерухомість з привабливою ставкою капіталізації.\n\n"
           "Ціна $3 000 000 відповідає ринковій вартості при базовому сценарії 6,5-річної окупності.\n\n"
           "Рекомендується провести комерційну due diligence щодо умов договору Варус "
           "та переговорів щодо пролонгації оренди землі.")
    T(s, rec,
      rx + Inches(0.15), ry + Inches(0.55), rw - Inches(0.25), rh - Inches(0.65),
      9, color=LABEL, wrap=True, v="t")

    # Ключові метрики під рекомендацією
    mets = [(f"$3M", "Ціна"), ("15,27%", "Cap Rate"), ("6,5 р.", "Окупність")]
    mw = rw / 3; my = ry + rh - Inches(0.60)
    R(s, rx, my - Inches(0.06), rw, Pt(1), GOLD_S)
    for i, (v, l) in enumerate(mets):
        mx = rx + i * mw
        T(s, v, mx, my, mw, Inches(0.30),
          12, bold=True, color=GOLD_T, align=PP_ALIGN.CENTER, v="ctr", wrap=False)
        T(s, l, mx, my + Inches(0.30), mw, Inches(0.26),
          8, color=MUTED, align=PP_ALIGN.CENTER, v="ctr", wrap=False)

    footer_bar(s, f"Ціна: $3{NB}000{NB}000  |  NOI: $458{NB}131/рік  |  Cap Rate: 15,27%  |  Окупність: 6,5 р.  |  IRR (10 р.): 20,49%")


# ═══════════════════════════════════════════════════════════════
# СЛАЙДИ 11-14 — ФОТОГАЛЕРЕЯ
# ═══════════════════════════════════════════════════════════════

def gallery_slide(prs, p1, cap1, p2, cap2, num, total=4):
    """2 фото в ряд на темному фоні."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    gold_vbar(s)

    # Мінімальний хедер (не bar, а просто тег)
    T(s, "ФОТОГАЛЕРЕЯ",
      Inches(0.28), Inches(0.06), Inches(6.0), Inches(0.38),
      10, bold=True, color=WHITE, v="ctr")
    T(s, f"{num}/{total}",
      Inches(8.5), Inches(0.06), Inches(1.3), Inches(0.38),
      9, color=GOLD_T, align=PP_ALIGN.RIGHT, v="ctr", wrap=False)
    T(s, "ТЦ «Бродвей»",
      Inches(6.0), Inches(0.06), Inches(2.3), Inches(0.38),
      8, color=MUTED, align=PP_ALIGN.RIGHT, v="ctr", wrap=False)

    # Розподіл фото
    PY = Inches(0.50); PH = H - PY - Inches(0.58)
    GAP = Inches(0.08)
    PW = (W - Inches(0.22) - GAP) / 2

    # Фото 1
    pic(s, p1, Inches(0.22), PY, PW, PH)
    # Фото 2
    pic(s, p2, Inches(0.22) + PW + GAP, PY, PW, PH)

    # Підписи
    cap_y = PY + PH + Inches(0.04)
    cap_h = H - cap_y - Inches(0.04)
    T(s, cap1, Inches(0.22), cap_y, PW, cap_h,
      9, color=MUTED, align=PP_ALIGN.CENTER, v="ctr", wrap=False)
    T(s, cap2, Inches(0.22) + PW + GAP, cap_y, PW, cap_h,
      9, color=MUTED, align=PP_ALIGN.CENTER, v="ctr", wrap=False)

    # Золота лінія між фото
    R(s, Inches(0.22) + PW, PY, GAP, PH, BG)
    R(s, Inches(0.22) + PW + GAP/2 - Pt(0.5), PY, Pt(1), PH, GOLD_S)


# ═══════════════════════════════════════════════════════════════
# СЛАЙД 15 — КОНТАКТИ
# ═══════════════════════════════════════════════════════════════

def slide15_contacts(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    gold_vbar(s)

    # Фонова напівпрозора панель
    R(s, Inches(0.22), 0, W - Inches(0.22), H, BG)

    # DE COSTA GROUP великим текстом по центру
    T(s, "DE COSTA GROUP",
      Inches(0.30), Inches(0.65), W - Inches(0.40), Inches(1.20),
      46, bold=True, color=GOLD_T, align=PP_ALIGN.CENTER, v="ctr")

    # Підзаголовок
    T(s, "Комерційна нерухомість України",
      Inches(0.30), Inches(1.78), W - Inches(0.40), Inches(0.40),
      14, color=LABEL, align=PP_ALIGN.CENTER, v="ctr")

    # Золота лінія
    R(s, Inches(1.50), Inches(2.28), W - Inches(3.00), Pt(1.5), GOLD_S)

    # Контактні дані (по центру)
    contacts = [
        ("✆",  "+38 (0xx) xxx-xx-xx"),
        ("✉",  "info@decosta.com.ua"),
        ("@",  "decosta.com.ua"),
        ("Pin", "Київ, Україна"),
    ]
    cy = Inches(2.50); cw = Inches(4.0); cx = (W - cw) / 2
    ch = Inches(0.52); cg = Inches(0.08)
    for icon, text in contacts:
        R(s, cx, cy, cw, ch, CARD)
        R(s, cx, cy, Inches(0.06), ch, GOLD_S)
        T(s, icon, cx + Inches(0.12), cy, Inches(0.45), ch,
          11, color=GOLD_T, align=PP_ALIGN.CENTER, v="ctr")
        T(s, text, cx + Inches(0.62), cy, cw - Inches(0.70), ch,
          11, color=WHITE, v="ctr", wrap=False)
        cy += ch + cg

    # Об'єкт що розглядається
    T(s, "ТЦ Бродвей  ·  вул. Таращанська, 191-А  ·  Біла Церква",
      Inches(0.30), H - Inches(0.55), W - Inches(0.40), Inches(0.30),
      8, color=MUTED, align=PP_ALIGN.CENTER, v="ctr", wrap=False)

    footer_bar(s, "Усі розрахунки виконані на підставі наданих документів та відкритих ринкових даних  |  Курс UAH/USD: 43,5")


# ═══════════════════════════════════════════════════════════════
# ЗБІРКА
# ═══════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # ── Основні слайди ─────────────────────────────────────────
    slide1_cover(prs)
    slide2_overview(prs)
    slide3_location(prs)
    slide4_technical(prs)
    slide5_legal(prs)
    slide6_tenants(prs)
    slide7_finance(prs)
    slide8_market(prs)
    slide9_swot(prs)
    slide10_conclusion(prs)

    # ── Фотогалерея ────────────────────────────────────────────
    gallery_slide(prs, P["main"],  "Головний фасад (вул. Таращанська)",
                       P["p1"],   "Фасад — вигляд збоку та паркінг", 1)
    gallery_slide(prs, P["p2"],   "Зовнішня аркада — галерея магазинів",
                       P["p5"],   "Торговий коридор — зона мала-ритейл", 2)
    gallery_slide(prs, P["p3"],   "Внутрішня торгова галерея",
                       P["p4"],   "Площа якірного орендаря (Варус, 1 751 м²)", 3)
    gallery_slide(prs, P["geo"],  "Супутникова карта локації (Заріччя)",
                       P["main"], "ТЦ Бродвей — зовнішній вигляд", 4)

    # ── Контакти ───────────────────────────────────────────────
    slide15_contacts(prs)

    out = "/sessions/zealous-beautiful-volta/broadway_full.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
