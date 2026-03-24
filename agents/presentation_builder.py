"""
Агент-построитель презентаций: генерирует PowerPoint или Gamma аутлайн.
Выходные файлы:
  - output/presentation.pptx (для format="pptx")
  - output/gamma_outline.txt (для format="gamma")
"""

from typing import Optional, Dict, Any
from pathlib import Path

from schemas import PropertyData, PipelineConfig


def build_presentation(
    property_data: PropertyData,
    metrics: Optional[Dict[str, Any]],
    config: PipelineConfig,
    output_dir: str = "output",
) -> str:
    """
    Генерирует презентацию в выбранном формате.

    Args:
        property_data: Извлечённые данные об объекте
        metrics: Финансовые метрики (или None для перспективных объектов)
        config: Конфигурация pipeline
        output_dir: Директория для выходных файлов

    Returns:
        Путь к созданному файлу
    """

    if config["pres_format"] == "pptx":
        return _build_pptx(property_data, metrics, config, output_dir)
    else:  # "gamma"
        return _build_gamma_outline(property_data, metrics, config, output_dir)


def _build_gamma_outline(
    property_data: PropertyData,
    metrics: Optional[Dict[str, Any]],
    config: PipelineConfig,
    output_dir: str,
) -> str:
    """Генерирует текстовый аутлайн для Gamma."""

    lines = []
    doc_type = config["doc_type"]
    object_category = config["object_category"]
    currency = config["currency"]
    currency_symbol = "$" if currency == "USD" else "₴"

    # Заголовок
    if object_category == "income":
        title = f"Інвестиційний аналіз: {property_data['property_name']}"
    else:
        title = f"Огляд об'єкту: {property_data['property_name']}"

    lines.append(title)
    lines.append(f"Тип об'єкту: {property_data['property_type']}, {property_data['city']}")
    lines.append("")

    # Слайды зависят от типа документа и категории объекта
    if doc_type == "teaser":
        lines.extend(
            _slides_teaser(property_data, metrics, object_category, currency_symbol)
        )
    elif doc_type == "memo":
        lines.extend(
            _slides_memo(property_data, metrics, object_category, currency_symbol)
        )
    else:  # "full"
        lines.extend(
            _slides_full(property_data, metrics, object_category, currency_symbol)
        )

    content = "\n".join(lines)

    # Сохраняем в файл
    output_path = Path(output_dir) / "gamma_outline.txt"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(content, encoding="utf-8")

    return str(output_path)


def _slides_teaser(
    property_data: PropertyData,
    metrics: Optional[Dict[str, Any]],
    object_category: str,
    currency_symbol: str,
) -> list:
    """Слайды для тизера (3-5 стр)."""
    lines = []

    lines.append("Слайд 1 — Обложка")
    lines.append(f"- Название: {property_data['property_name']}")
    lines.append(f"- Тип: {property_data['property_type']}")
    lines.append(f"- Локация: {property_data['city']}")
    lines.append("")

    lines.append("Слайд 2 — Ключевые характеристики")
    lines.append(f"- GBA: {property_data['gba']:,.0f} м²")
    if property_data.get("gla"):
        lines.append(f"- GLA: {property_data['gla']:,.0f} м²")
    if property_data.get("year_built"):
        lines.append(f"- Год постройки: {property_data['year_built']}")
    if property_data.get("condition"):
        lines.append(f"- Состояние: {property_data['condition']}")
    if property_data.get("infrastructure"):
        lines.append(f"- Локация: {property_data['infrastructure']}")
    lines.append("")

    if object_category == "income" and metrics:
        lines.append("Слайд 3 — Финансовый snapshot")
        lines.append(f"- Стоимость: {metrics['value']:,.0f} {currency_symbol}")
        if metrics.get("noi"):
            lines.append(f"- NOI: {metrics['noi']:,.0f} {currency_symbol}/год")
        if metrics.get("cap_rate"):
            lines.append(f"- Cap Rate: {metrics['cap_rate']:.2f}%")
        lines.append("")

    return lines


def _slides_memo(
    property_data: PropertyData,
    metrics: Optional[Dict[str, Any]],
    object_category: str,
    currency_symbol: str,
) -> list:
    """Слайды для меморандума (5-10 стр)."""
    lines = []

    lines.append("Слайд 1 — Обложка")
    lines.append(f"- {property_data['property_name']}")
    lines.append(f"- {property_data['property_type'].upper()}")
    lines.append("")

    lines.append("Слайд 2 — Резюме")
    lines.append(f"- {property_data.get('description', 'Описание недоступно')}")
    lines.append("")

    lines.append("Слайд 3 — Характеристики")
    lines.append(f"- GBA: {property_data['gba']:,.0f} м²")
    if property_data.get("gla"):
        lines.append(f"- GLA: {property_data['gla']:,.0f} м²")
    if property_data.get("year_built"):
        lines.append(f"- Год: {property_data['year_built']}")
    if property_data.get("infrastructure"):
        lines.append(f"- Инфраструктура: {property_data['infrastructure']}")
    lines.append("")

    if object_category == "income" and metrics:
        lines.append("Слайд 4 — Финансовые показатели")
        lines.append(f"- Стоимость: {metrics['value']:,.0f} {currency_symbol}")
        if metrics.get("noi"):
            lines.append(f"- NOI: {metrics['noi']:,.0f} {currency_symbol}/год")
        if metrics.get("cap_rate"):
            lines.append(f"- Cap Rate: {metrics['cap_rate']:.2f}%")
        if metrics.get("payback_years"):
            lines.append(f"- Окупаемость: {metrics['payback_years']:.1f} лет")
        lines.append("")

        if metrics.get("cap_rate_scenarios"):
            lines.append("Слайд 5 — Сценарии оценки")
            for scenario in metrics["cap_rate_scenarios"]:
                # Сценарии могут быть dict или список
                if isinstance(scenario, dict):
                    cap_rate = scenario.get("cap_rate_pct")
                    value = scenario.get("property_value")
                elif isinstance(scenario, (list, tuple)) and len(scenario) >= 2:
                    cap_rate = scenario[0]
                    value = scenario[1]
                else:
                    continue

                if cap_rate and value:
                    lines.append(
                        f"- Cap Rate {cap_rate:.1f}%: {value:,.0f} {currency_symbol}"
                    )
            lines.append("")

    return lines


def _slides_full(
    property_data: PropertyData,
    metrics: Optional[Dict[str, Any]],
    object_category: str,
    currency_symbol: str,
) -> list:
    """Слайды для полной инвестиционной презентации (10+ стр)."""
    lines = []

    lines.append("Слайд 1 — Обложка")
    lines.append(f"- {property_data['property_name']}")
    lines.append(f"- Инвестиционный анализ")
    lines.append("")

    lines.append("Слайд 2 — Резюме")
    lines.append(f"- {property_data.get('description', 'Описание')}")
    lines.append("")

    lines.append("Слайд 3 — Обзор объекта")
    lines.append(f"- Адрес: {property_data.get('address', 'N/A')}")
    lines.append(f"- GBA: {property_data['gba']:,.0f} м²")
    lines.append(f"- Тип: {property_data['property_type']}")
    if property_data.get("year_built"):
        lines.append(f"- Год постройки: {property_data['year_built']}")
    lines.append("")

    if object_category == "income" and metrics:
        lines.append("Слайд 4 — Финансовые показатели")
        lines.append(f"- Стоимость: {metrics['value']:,.0f} {currency_symbol}")
        if metrics.get("noi"):
            lines.append(f"- NOI: {metrics['noi']:,.0f} {currency_symbol}")
        if metrics.get("cap_rate"):
            lines.append(f"- Cap Rate: {metrics['cap_rate']:.2f}%")
        lines.append("")

        lines.append("Слайд 5 — Детальная финансовая модель")
        lines.append(f"- Валовой доход: {metrics.get('pgi', 'N/A')}")
        lines.append(f"- Расходы: {metrics.get('opex', 'N/A')}")
        lines.append(f"- Чистый доход: {metrics.get('noi', 'N/A')}")
        lines.append("")

        if metrics.get("cap_rate_scenarios"):
            lines.append("Слайд 6 — Сценарии (Cap Rate)")
            for scenario in metrics["cap_rate_scenarios"]:
                # Сценарии могут быть dict или список
                if isinstance(scenario, dict):
                    cap_rate = scenario.get("cap_rate_pct")
                    value = scenario.get("property_value")
                elif isinstance(scenario, (list, tuple)) and len(scenario) >= 2:
                    cap_rate = scenario[0]
                    value = scenario[1]
                else:
                    continue

                if cap_rate and value:
                    lines.append(f"- {cap_rate:.1f}%: {value:,.0f} {currency_symbol}")
            lines.append("")

        lines.append("Слайд 7 — Анализ рисков")
        lines.append("- Рыночные риски")
        lines.append("- Операционные риски")
        lines.append("- Финансовые риски")
        lines.append("")

        lines.append("Слайд 8 — Выводы и рекомендации")
        lines.append("- Потенциал доходности")
        lines.append("- Инвестиционный рейтинг")
        lines.append("")

    else:
        lines.append("Слайд 4 — Рынок и анализ")
        lines.append("- Анализ рыночного спроса")
        lines.append("- Конкурентная среда")
        lines.append("- Возможности использования")
        lines.append("")

        lines.append("Слайд 5 — Техническое описание")
        lines.append(f"- Состояние: {property_data.get('condition', 'N/A')}")
        if property_data.get("infrastructure"):
            lines.append(f"- Инфраструктура: {property_data['infrastructure']}")
        lines.append("")

        lines.append("Слайд 6 — Выводы")
        lines.append("- Возможности развития")
        lines.append("- Рекомендации")
        lines.append("")

    return lines


def _build_pptx(
    property_data: PropertyData,
    metrics: Optional[Dict[str, Any]],
    config: PipelineConfig,
    output_dir: str,
) -> str:
    """Генерирует PowerPoint файл."""

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError(
            "Для создания PowerPoint требуется установить: pip3 install python-pptx"
        )

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Добавляем слайды в зависимости от типа и категории
    doc_type = config["doc_type"]
    object_category = config["object_category"]
    currency = config["currency"]
    currency_symbol = "$" if currency == "USD" else "₴"

    if doc_type == "teaser":
        _add_teaser_slides(prs, property_data, metrics, object_category, currency_symbol)
    elif doc_type == "memo":
        _add_memo_slides(prs, property_data, metrics, object_category, currency_symbol)
    else:  # "full"
        _add_full_slides(prs, property_data, metrics, object_category, currency_symbol)

    # Сохраняем
    output_path = Path(output_dir) / "presentation.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    return str(output_path)


def _add_teaser_slides(prs, data, metrics, category, curr_sym):
    """Добавляет слайды для тизера."""
    # Слайд 1 — Обложка
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой лейаут
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = data["property_name"]
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True

    # Слайд 2 — Характеристики
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Характеристики"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = f"GBA: {data['gba']:,.0f} м²"
    if data.get("gla"):
        tf.add_paragraph().text = f"GLA: {data['gla']:,.0f} м²"
    if data.get("year_built"):
        tf.add_paragraph().text = f"Год постройки: {data['year_built']}"


def _add_memo_slides(prs, data, metrics, category, curr_sym):
    """Добавляет слайды для меморандума."""
    # Слайд 1 — Обложка
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = data["property_name"]
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True


def _add_full_slides(prs, data, metrics, category, curr_sym):
    """Добавляет слайды для полной презентации."""
    # Слайд 1 — Обложка
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = data["property_name"]
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    stf = subtitle_box.text_frame
    stf.text = "Инвестиционный анализ"
    stf.paragraphs[0].font.size = Pt(28)


if __name__ == "__main__":
    from schemas import PropertyData

    test_data = PropertyData(
        property_name="Склад Тест",
        property_type="склад",
        address="вул. Тестова, 1",
        city="Київ",
        description="Тестовий складський комплекс",
        gba=5000,
        gla=4500,
        value=1200000,
        value_currency="USD",
        rent_rate=6,
        rent_rate_currency="USD",
        vacancy_rate=0.05,
        opex=50000,
        opex_currency="USD",
        rent_growth_rate=0.03,
        exit_cap_rate=0.09,
        hold_period=10,
        year_built=2015,
        condition="задовільний",
        land_area=None,
        infrastructure="біля МКАД",
        source_file="test.txt",
        extraction_confidence="high",
        extraction_notes="",
    )

    test_metrics = {
        "value": 1200000,
        "noi": 86400,
        "cap_rate": 7.2,
        "payback_years": 13.9,
    }

    config = PipelineConfig(
        object_category="income",
        doc_type="memo",
        pres_format="gamma",
        exchange_rate=41.5,
        currency="USD",
    )

    output_file = build_presentation(test_data, test_metrics, config, output_dir="/tmp")
    print(f"Створен файл: {output_file}")
